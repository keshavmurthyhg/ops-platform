"""
ppt_creator.py — Clean landscape RCA PowerPoint
Slide 1:  Incident info with proper tables (hyperlinks for INC/Azure/PTC)
Slides 2-4: Problem / Root Cause / Resolution — text + images on ONE slide
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

from common.logger import setup_logger
from common.utils.links import parse_ptc_cases as _parse_ptc_cases
logger = setup_logger("ppt_creator")


def _format_date(val):
    """Format date to DD-MMM-YYYY, matching word_renderer.format_date()."""
    if not val or str(val).strip() in ("", "-", "None", "nan", "NaT", "nat"):
        return "-"
    try:
        from datetime import datetime
        # Strip timestamp if present (e.g. "2026-01-22 08:57:07" → "2026-01-22")
        date_part = str(val).strip().split()[0]
        return datetime.strptime(date_part, "%Y-%m-%d").strftime("%d-%b-%Y")
    except Exception:
        return str(val).strip()

# ── Dimensions ────────────────────────────────────────────────────────────────
SLIDE_W   = Inches(13.33)
SLIDE_H   = Inches(7.5)
M         = Inches(0.45)          # left/right margin
CONTENT_W = SLIDE_W - 2 * M

# ── Colour palette ────────────────────────────────────────────────────────────
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK    = RGBColor(0x1E, 0x29, 0x3B)
C_GREY_HD  = RGBColor(0xF1, 0xF5, 0xF9)   # header cell fill
C_GREY_ROW = RGBColor(0xF8, 0xFA, 0xFC)   # alt row fill
C_BORDER   = RGBColor(0xCB, 0xD5, 0xE1)
C_LINK     = RGBColor(0x2F, 0x80, 0xED)
C_MUTED    = RGBColor(0x64, 0x74, 0x8B)   # secondary text

C_PROBLEM  = RGBColor(0xDC, 0x26, 0x26)
C_ROOT     = RGBColor(0xD9, 0x77, 0x06)
C_RESOLVE  = RGBColor(0x16, 0xA3, 0x4A)
C_TITLE_BG = RGBColor(0x1E, 0x29, 0x3B)   # title band (slide 1)

# ── URL base configuration ────────────────────────────────────────────────────
# These must match what common.utils.links.apply_word_link uses in Word/PDF.
# Update BASE_AZURE_URL and BASE_PTC_URL if your organisation uses different paths.
BASE_SNOW_URL  = "https://volvoitsm.service-now.com/nav_to.do?uri=incident.do?sysparm_query=number="
BASE_AZURE_URL = "https://dev.azure.com/VolvoCarsGroup/VCE/_workitems/edit/"
BASE_PTC_URL   = "https://support.ptc.com/appserver/cs/view/case.jsp?n="


# ══════════════════════════════════════════════════════════════════════════════
# Primitives
# ══════════════════════════════════════════════════════════════════════════════

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rgb_hex(rgb):
    return '%02X%02X%02X' % (rgb[0], rgb[1], rgb[2])


def _set_fill(shape, rgb):
    sp = shape._element.spPr
    for ch in list(sp):
        tag = ch.tag.split("}")[-1]
        if tag in ('solidFill','gradFill','noFill','blipFill','pattFill'):
            sp.remove(ch)
    sf  = etree.SubElement(sp, qn('a:solidFill'))
    clr = etree.SubElement(sf,  qn('a:srgbClr'))
    clr.set('val', _rgb_hex(rgb))


def _rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(0.5)):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        _set_fill(sh, fill)
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border
        sh.line.width = bw
    else:
        sh.line.fill.background()
    return sh


def _tb(slide, text, l, t, w, h,
        sz=Pt(10), bold=False, color=None, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    color = color or C_BLACK
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text          = str(text)
    r.font.size     = sz
    r.font.bold     = bold
    r.font.italic   = italic
    r.font.color.rgb = color
    return txb


def _link_tb(slide, text, url, l, t, w, h, sz=Pt(10)):
    """Textbox with a clickable hyperlink."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    r   = p.add_run()
    r.text           = str(text)
    r.font.size      = sz
    r.font.color.rgb = C_BLACK
    r.font.underline = False
    if url:
        rId = txb.part.relate_to(
            url,
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
            is_external=True
        )
        rPr = r._r.get_or_add_rPr()
        hl  = etree.SubElement(rPr, qn('a:hlinkClick'))
        hl.set(qn('r:id'), rId)
    return txb


def _multi_link_tb(slide, cases, base_url, l, t, w, h, sz=Pt(9.5)):
    """Textbox with multiple comma-separated hyperlinks (for PTC multi-case).
    cases: list of (display_text, numeric_id) tuples.
    """
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    for i, (display, num_id) in enumerate(cases):
        if i > 0:
            sep = p.add_run()
            sep.text       = ", "
            sep.font.size  = sz
            sep.font.color.rgb = C_BLACK
        r = p.add_run()
        r.text           = display
        r.font.size      = sz
        r.font.color.rgb = C_BLACK
        r.font.underline = False
        url = base_url + num_id
        rId = txb.part.relate_to(
            url,
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
            is_external=True
        )
        rPr = r._r.get_or_add_rPr()
        hl  = etree.SubElement(rPr, qn('a:hlinkClick'))
        hl.set(qn('r:id'), rId)
    return txb


def _accent_bar(slide, rgb, thick=Pt(4)):
    """Thin coloured bar at very top of slide."""
    _rect(slide, 0, 0, SLIDE_W, thick, fill=rgb)


def _section_head(slide, emoji, title, accent, top):
    """Grey heading band with small coloured circle + title text."""
    band_h = Inches(0.52)
    _rect(slide, 0, top, SLIDE_W, band_h, fill=C_GREY_HD)
    # Small filled circle (10pt bullet in accent colour)
    _tb(slide, "●", M, top + Inches(0.1), Inches(0.3), band_h - Inches(0.1),
        sz=Pt(10), color=accent)
    _tb(slide, title, M + Inches(0.3), top + Inches(0.06),
        CONTENT_W - Inches(0.3), band_h - Inches(0.1),
        sz=Pt(15), bold=True, color=C_BLACK)
    return top + band_h


# ══════════════════════════════════════════════════════════════════════════════
# Data helpers
# ══════════════════════════════════════════════════════════════════════════════

def _g(d, *keys):
    """Get first non-empty value from dict trying multiple key variants."""
    for k in keys:
        for variant in [k, k.upper(), k.lower(), k.title(),
                        k.replace("_"," "), k.replace(" ","_")]:
            v = d.get(variant)
            if v and str(v).strip() not in ("-", "", "None", "nan", "NaT", "nat"):
                return str(v).strip()
    return "-"


def _cell(slide, text, l, t, w, h, bold=False, fill=None,
          link_url=None, sz=Pt(9.5), color=None, pad=Inches(0.07)):
    """Single table cell = rect + textbox."""
    _rect(slide, l, t, w, h, fill=fill or C_WHITE, border=C_BORDER, bw=Pt(0.5))
    if link_url and text != "-":
        _link_tb(slide, text, link_url, l+pad, t+pad, w-2*pad, h-2*pad, sz=sz)
    else:
        _tb(slide, text, l+pad, t+pad, w-2*pad, h-2*pad,
            sz=sz, bold=bold, color=color or C_BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Incident Info
# ══════════════════════════════════════════════════════════════════════════════

def _slide_incident(prs, data, inc_num):
    slide = _blank(prs)
    logger.info("Building incident slide: %s", inc_num)

    # ── Title band ────────────────────────────────────────────────────────────
    TITLE_H = Inches(0.48)
    _rect(slide, 0, 0, SLIDE_W, TITLE_H, fill=C_TITLE_BG)
    _tb(slide, f"INCIDENT REPORT  —  {inc_num}", M, Inches(0.06),
        CONTENT_W, TITLE_H - Inches(0.06),
        sz=Pt(16), bold=True, color=C_WHITE)

    top = TITLE_H + Inches(0.1)

    # ── Metadata table ────────────────────────────────────────────────────────
    # Exact keys from data_mapper.py:
    # number, created_by, azure_bug, created_date, ptc_case, assigned_to,
    # priority, resolved_date
    # Metadata rows — exactly matching word_renderer column layout (1.0+2.0+1.5+2.0 = 6.5")
    # Left pair total = LBL_W + VAL_W; right pair total = LBL_W + VAL_W
    # Both pairs side by side span CONTENT_W
    # Word column widths: 1.0 | 2.0 | 1.5 | 2.0  (total 6.5")
    # We scale proportionally to our CONTENT_W (12.43")
    # Ratios: label=1.0/6.5, value=2.0/6.5, label2=1.5/6.5, value2=2.0/6.5
    TOTAL_W   = CONTENT_W
    LBL_W     = TOTAL_W * (1.0 / 6.5)
    VAL_W     = TOTAL_W * (2.0 / 6.5)
    LBL_W2    = TOTAL_W * (1.5 / 6.5)
    VAL_W2    = TOTAL_W * (2.0 / 6.5)
    ROW_H     = Inches(0.3)

    META_ROWS = [
        # (lbl_left, key_left, is_link, lbl_right, key_right, is_date_right)
        # is_link=True → will generate a hyperlink using lbl1 to determine URL base
        ("INCIDENT",      "number",       True,  "CREATED BY",    "created_by",    False),
        ("AZURE BUG",     "azure_bug",    True,  "CREATED DATE",  "created_date",  True),
        ("PTC CASE",      "ptc_case",     True,  "ASSIGNED TO",   "assigned_to",   False),
        ("PRIORITY",      "priority",     False, "RESOLVED DATE", "resolved_date", True),
    ]

    for ri, (lbl1, key1, link1, lbl2, key2, date2) in enumerate(META_ROWS):
        val1 = _g(data, key1, lbl1)
        val2_raw = _g(data, key2, lbl2)
        val2 = _format_date(val2_raw) if date2 else val2_raw
        y   = top + ri * ROW_H
        alt = C_GREY_ROW if ri % 2 == 0 else C_WHITE

        # Left pair
        x0 = M
        _cell(slide, lbl1, x0,           y, LBL_W,  ROW_H, bold=True, fill=C_GREY_HD)

        # PTC CASE gets special multi-link treatment
        if "ptc" in lbl1.lower():
            ptc_cases = _parse_ptc_cases(val1 if val1 != "-" else data.get("ptc_case", ""))
            _rect(slide, x0 + LBL_W, y, VAL_W, ROW_H, fill=alt, border=C_BORDER, bw=Pt(0.5))
            pad = Inches(0.07)
            if ptc_cases:
                _multi_link_tb(slide, ptc_cases, BASE_PTC_URL,
                               x0 + LBL_W + pad, y + pad,
                               VAL_W - 2*pad, ROW_H - 2*pad, sz=Pt(9.5))
            else:
                _tb(slide, "-", x0 + LBL_W + pad, y + pad,
                    VAL_W - 2*pad, ROW_H - 2*pad, sz=Pt(9.5))
        else:
            # Build hyperlink URL matching Word/PDF behaviour
            if link1 and val1 != "-":
                lbl_lower = lbl1.lower()
                if "incident" in lbl_lower:
                    url1 = BASE_SNOW_URL + val1
                elif "azure" in lbl_lower:
                    url1 = BASE_AZURE_URL + val1
                else:
                    url1 = None
            else:
                url1 = None
            _cell(slide, val1, x0 + LBL_W, y, VAL_W, ROW_H, fill=alt, link_url=url1)

        # Right pair
        x1 = M + LBL_W + VAL_W
        _cell(slide, lbl2, x1,            y, LBL_W2, ROW_H, bold=True, fill=C_GREY_HD)
        _cell(slide, val2,  x1 + LBL_W2,  y, VAL_W2, ROW_H, fill=alt)

    top = top + len(META_ROWS) * ROW_H + Inches(0.14)

    # ── Short Description / Description table ─────────────────────────────────
    short = _g(data, "short_description", "SHORT DESCRIPTION", "Short Description")
    desc  = _g(data, "description", "DESCRIPTION", "Description")
    # Match Word: short desc col = 1.0+2.0=3.0", desc col = 1.5+2.0=3.5"
    SD_W   = (LBL_W  + VAL_W)    # 3.0" proportional
    DC_W   = (LBL_W2 + VAL_W2)   # 3.5" proportional
    DROW_H = Inches(0.76)
    HDR_H  = Inches(0.27)

    _cell(slide, "SHORT DESCRIPTION", M,        top,        SD_W, HDR_H, bold=True, fill=C_GREY_HD, sz=Pt(9))
    _cell(slide, "DESCRIPTION",       M + SD_W, top,        DC_W, HDR_H, bold=True, fill=C_GREY_HD, sz=Pt(9))
    _cell(slide, short[:180],         M,        top+HDR_H, SD_W, DROW_H, sz=Pt(9))
    _cell(slide, desc[:260],          M + SD_W, top+HDR_H, DC_W, DROW_H, sz=Pt(9))

    top = top + HDR_H + DROW_H + Inches(0.14)

    # ── RCA Summary table (3-row: Problem / Root Cause / Resolution) ──────────
    problem_text   = _g(data, "problem",   "problem_statement", "PROBLEM STATEMENT")
    rootcause_text = _g(data, "analysis",  "root_cause",        "ROOT CAUSE")
    resolution_txt = _g(data, "resolution","RESOLUTION")

    RCA_ROWS = [
        ("PROBLEM STATEMENT", problem_text,   C_PROBLEM),
        ("ROOT CAUSE",        rootcause_text, C_ROOT),
        ("RESOLUTION",        resolution_txt, C_RESOLVE),
    ]
    RCA_LBL_W = Inches(2.1)
    RCA_VAL_W = CONTENT_W - RCA_LBL_W    # same total width as meta/desc tables
    RCA_ROW_H = Inches(0.58)

    for ri, (lbl, val, acc) in enumerate(RCA_ROWS):
        y    = top + ri * RCA_ROW_H
        fill = C_GREY_ROW if ri % 2 == 0 else C_WHITE
        stripe_w = Inches(0.05)
        _rect(slide, M, y, stripe_w, RCA_ROW_H, fill=acc)
        _cell(slide, lbl, M + stripe_w, y, RCA_LBL_W - stripe_w, RCA_ROW_H,
              bold=True, fill=C_GREY_HD, sz=Pt(9))
        _cell(slide, val[:400], M + RCA_LBL_W, y, RCA_VAL_W, RCA_ROW_H,
              fill=fill, sz=Pt(9))

    top = top + len(RCA_ROWS) * RCA_ROW_H + Inches(0.1)

    # ── References table — styled to match metadata/description/RCA tables ──
    refs = (data or {}).get("references") or []
    if refs and (SLIDE_H - top) > Inches(0.45):
        avail_h = SLIDE_H - top - Inches(0.05)

        # Header row height + data rows
        HDR_H  = Inches(0.27)
        # Distribute remaining space equally across ref rows (min 0.28", max 0.38")
        max_rows = int((avail_h - HDR_H) / Inches(0.28))
        visible  = refs[:max(1, max_rows)]
        ROW_H    = min(Inches(0.38), (avail_h - HDR_H) / max(len(visible), 1))

        # Column widths proportional to slide content width
        # Ref=30%, Env=12%, Link&Context=58%
        REF_W = CONTENT_W * 0.30
        ENV_W = CONTENT_W * 0.12
        LNK_W = CONTENT_W * 0.58

        REF_COLORS = {
            "azure_user_story": RGBColor(0x25, 0x63, 0xEB),
            "ptc_article":      RGBColor(0x93, 0x33, 0xEA),
        }

        # Section header band — same purple style as other section heads
        _rect(slide, 0, top, SLIDE_W, Inches(0.22), fill=RGBColor(0xED, 0xE9, 0xFE))
        _tb(slide, "🔗  REFERENCES", M, top + Inches(0.02),
            CONTENT_W, Inches(0.20), sz=Pt(9), bold=True,
            color=RGBColor(0x5B, 0x21, 0xB6))
        top += Inches(0.22)

        # Header row — same fill/font as metadata & RCA label cells (C_GREY_HD, bold, Pt(9), C_BLACK)
        x = M
        _cell(slide, "REFERENCE",    x,           top, REF_W, HDR_H, bold=True, fill=C_GREY_HD, sz=Pt(9), color=C_BLACK)
        _cell(slide, "ENVIRONMENT",  x + REF_W,   top, ENV_W, HDR_H, bold=True, fill=C_GREY_HD, sz=Pt(9), color=C_BLACK)
        _cell(slide, "LINK & CONTEXT", x+REF_W+ENV_W, top, LNK_W, HDR_H, bold=True, fill=C_GREY_HD, sz=Pt(9), color=C_BLACK)
        top += HDR_H

        # Data rows
        for ri, ref in enumerate(visible):
            if top + ROW_H > SLIDE_H - Inches(0.05):
                break
            row_fill  = C_GREY_ROW if ri % 2 == 0 else C_WHITE
            ref_color = REF_COLORS.get(ref.get("type", ""), C_MUTED)
            env_val   = ref.get("environment") or "-"
            ctx       = ref.get("context", "")[:90]

            # Col 0: Reference label — coloured text matching type, same Pt(9) as other tables
            _rect(slide, x, top, REF_W, ROW_H, fill=row_fill, border=C_BORDER, bw=Pt(0.4))
            _tb(slide, ref.get("label", ""), x + Inches(0.07), top + Inches(0.04),
                REF_W - Inches(0.1), ROW_H - Inches(0.06),
                sz=Pt(8.5), bold=True, color=ref_color)

            # Col 1: Environment — plain Pt(9), C_BLACK
            _cell(slide, env_val, x + REF_W, top, ENV_W, ROW_H, fill=row_fill, sz=Pt(8.5), color=C_BLACK)

            # Col 2: Clickable URL + context
            _rect(slide, x+REF_W+ENV_W, top, LNK_W, ROW_H, fill=row_fill, border=C_BORDER, bw=Pt(0.4))
            pad = Inches(0.07)
            url = ref.get("url", "")
            if url:
                _link_tb(slide, url, url,
                         x+REF_W+ENV_W+pad, top+Inches(0.03),
                         LNK_W-2*pad, Inches(0.18), sz=Pt(8))
            if ctx:
                ctx_top = top + Inches(0.20)
                if ctx_top + Inches(0.15) < top + ROW_H:
                    _tb(slide, ctx, x+REF_W+ENV_W+pad, ctx_top,
                        LNK_W-2*pad, ROW_H-Inches(0.22),
                        sz=Pt(7.5), color=C_MUTED)

            top += ROW_H

    logger.info("Incident slide built")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# RCA SECTION SLIDES — text + images on SAME slide
# ══════════════════════════════════════════════════════════════════════════════

def _slides_rca(prs, zone_key, zone_title, accent, text, images):
    emoji_map = {"problem": "Problem", "rootcause": "Root Cause", "resolution": "Resolution"}
    logger.info("Building %s — %d chars, %d images", zone_key, len(text or ""), len(images))

    # Process in chunks of 3 images per slide (first slide also has text)
    chunks = [images[i:i+3] for i in range(0, max(1, len(images)), 3)] if images else [[]]

    for ci, chunk in enumerate(chunks):
        slide = _blank(prs)

        # 4pt accent bar
        _accent_bar(slide, accent)
        content_top = _section_head(slide, None, zone_title, accent, top=Pt(4))
        content_top += Inches(0.08)

        if ci == 0:
            # Text area
            body = (text or "").strip() or "No analysis text provided."
            text_h = Inches(2.1) if chunk else (SLIDE_H - content_top - Inches(0.2))

            _rect(slide, M, content_top, CONTENT_W, text_h,
                  fill=C_GREY_ROW, border=C_BORDER, bw=Pt(0.5))
            _tb(slide, body,
                M+Inches(0.12), content_top+Inches(0.1),
                CONTENT_W-Inches(0.24), text_h-Inches(0.18),
                sz=Pt(11), color=C_BLACK, wrap=True)
            content_top += text_h + Inches(0.1)

        # Images in horizontal strip
        if chunk:
            avail_h = SLIDE_H - content_top - Inches(0.12)
            avail_w = CONTENT_W
            n       = len(chunk)
            img_w   = (avail_w - Inches(0.08) * (n - 1)) / n

            for i, img_path in enumerate(chunk):
                if not os.path.exists(img_path):
                    logger.warning("Image not found: %s", img_path)
                    continue
                try:
                    x   = M + i * (img_w + Inches(0.08))
                    pic = slide.shapes.add_picture(img_path, x, content_top, width=img_w)
                    # Constrain height
                    if pic.height > avail_h:
                        ratio      = avail_h / pic.height
                        pic.height = int(pic.height * ratio)
                        pic.width  = int(pic.width  * ratio)
                        pic.left   = int(M + i * (img_w + Inches(0.08)))
                    # Centre vertically
                    pic.top = int(content_top + (avail_h - pic.height) / 2)
                    logger.info("  img %d/%d placed: %s", i+1, n, os.path.basename(img_path))
                except Exception as e:
                    logger.error("  img error %s: %s", img_path, e)

        if ci > 0:
            # Page label for overflow slides
            _tb(slide, f"(continued — images {ci*3+1}–{min(ci*3+3, len(images))} of {len(images)})",
                M, SLIDE_H - Inches(0.3), CONTENT_W, Inches(0.25),
                sz=Pt(8), color=C_MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN ENTRY
# ══════════════════════════════════════════════════════════════════════════════

def create_rca_pptx(incident_data, incident_number, rca_text, rca_image_paths, output_path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    logger.info("=" * 60)
    logger.info("RCA PPTX CREATION")
    logger.info("  Incident : %s", incident_number)
    logger.info("  Output   : %s", output_path)
    # Log what data keys are available (helps debug missing fields)
    logger.info("  Data keys: %s", list((incident_data or {}).keys())[:15])

    _slide_incident(prs, incident_data or {}, incident_number or "")

    SECTIONS = [
        ("problem",    "Problem Statement", C_PROBLEM),
        ("rootcause",  "Root Cause",        C_ROOT),
        ("resolution", "Resolution",        C_RESOLVE),
    ]
    for zone_key, zone_title, color in SECTIONS:
        text   = (rca_text or {}).get(zone_key, "")
        images = (rca_image_paths or {}).get(zone_key, [])
        logger.info("  %-12s: %d img, %d chars", zone_key, len(images), len(text))
        _slides_rca(prs, zone_key, zone_title, color, text, images)

    prs.save(output_path)
    kb = os.path.getsize(output_path) // 1024
    logger.info("PPTX SAVED: %s (%d KB)", output_path, kb)
    logger.info("=" * 60)
    return output_path
