import os
import tempfile
import subprocess
import shutil

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pdf2image import convert_from_path

from common.logger import setup_logger

logger = setup_logger("ppt_slide_renderer")

LIBREOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"
POPPLER_PATH     = r"C:\poppler\Library\bin"

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

_STRUCTURAL_TAGS = {
    f"{{{NS_P}}}nvGrpSpPr",
    f"{{{NS_P}}}grpSpPr",
}

# Human-readable shape type labels for logging
_SHAPE_TAG_LABELS = {
    f"{{{NS_P}}}sp":     "shape",
    f"{{{NS_P}}}pic":    "picture",
    f"{{{NS_P}}}cxnSp":  "connector",
    f"{{{NS_P}}}grpSp":  "group",
}


# ------------------------------------------------------------------
# Shape classification helpers
# ------------------------------------------------------------------

def _has_content_text(el):
    for t in el.iter(f"{{{NS_A}}}t"):
        if t.text and t.text.strip():
            return True
    return False


def _has_content_image(el):
    if el.tag == f"{{{NS_P}}}pic":
        return True
    if next(el.iter(f"{{{NS_P}}}pic"), None) is not None:
        return True
    if next(el.iter(f"{{{NS_A}}}blipFill"), None) is not None:
        return True
    return False


def _shape_label(el):
    """Return a short human-readable label for a shape element."""
    cNvPr = el.find(f".//{{{NS_P}}}cNvPr")
    name  = cNvPr.get("name", "") if cNvPr is not None else ""
    kind  = _SHAPE_TAG_LABELS.get(el.tag, el.tag.split("}")[1])
    return f"{kind}('{name}')" if name else kind


def _should_keep_in_master_or_layout(el):
    """
    Whitelist for master/layout spTree elements.
    Background decorative shapes (swooshes, chrome) live ONLY in master/layout
    — never in the slide spTree. So we strip them here but never on slides.
    """
    if el.tag in _STRUCTURAL_TAGS:
        return True
    if _has_content_image(el):
        return True
    if _has_content_text(el):
        return True
    if el.find(f".//{{{NS_P}}}ph") is not None:
        return True
    return False


# ------------------------------------------------------------------
# White background injection
# ------------------------------------------------------------------

def _make_white_bg():
    bg   = etree.Element(f"{{{NS_P}}}bg")
    bgPr = etree.SubElement(bg, f"{{{NS_P}}}bgPr")
    sf   = etree.SubElement(bgPr, f"{{{NS_A}}}solidFill")
    clr  = etree.SubElement(sf,   f"{{{NS_A}}}srgbClr")
    clr.set("val", "FFFFFF")
    return bg


def _set_white_bg(xml_el, label=""):
    cSld = xml_el.find(f"{{{NS_P}}}cSld")
    if cSld is None:
        return
    for existing in cSld.findall(f"{{{NS_P}}}bg"):
        cSld.remove(existing)
    spTree    = cSld.find(f"{{{NS_P}}}spTree")
    insert_at = list(cSld).index(spTree) if spTree is not None else 0
    cSld.insert(insert_at, _make_white_bg())
    logger.debug("%s: white background applied (bgRef removed)", label)


def _strip_master_layout_chrome(xml_el, label=""):
    """
    Remove decorative shapes from a master or layout spTree and log
    each removed shape by name/type.
    """
    cSld   = xml_el.find(f"{{{NS_P}}}cSld")
    if cSld is None:
        return
    spTree = cSld.find(f"{{{NS_P}}}spTree")
    if spTree is None:
        return

    stripped  = []
    kept_dec  = []

    for child in list(spTree):
        if _should_keep_in_master_or_layout(child):
            kept_dec.append(_shape_label(child))
        else:
            stripped.append(_shape_label(child))
            spTree.remove(child)

    if stripped:
        logger.info(
            "%s | stripped %d chrome shape(s): %s",
            label, len(stripped), ", ".join(stripped)
        )
    else:
        logger.debug("%s | no chrome shapes found", label)


# ------------------------------------------------------------------
# Slide dimension helper
# ------------------------------------------------------------------

def _slide_wh(slide):
    try:
        prs = slide.part.package.presentation_part.presentation
        return int(prs.slide_width), int(prs.slide_height)
    except Exception:
        return 9144000, 6858000


# ------------------------------------------------------------------
# Slide content summary for logging
# ------------------------------------------------------------------

def _slide_shape_summary(slide):
    """Return a dict counting shape types on a slide for logging."""
    counts = {
        "pictures":   0,
        "text_boxes": 0,
        "connectors": 0,
        "shapes":     0,
        "groups":     0,
        "total":      0,
    }
    NS_P_ = NS_P
    spTree = slide._element.find(f".//{{{NS_P_}}}spTree")
    if spTree is None:
        return counts
    for child in spTree:
        tag = child.tag
        if tag in _STRUCTURAL_TAGS:
            continue
        counts["total"] += 1
        if tag == f"{{{NS_P_}}}pic" or _has_content_image(child):
            counts["pictures"] += 1
        elif tag == f"{{{NS_P_}}}cxnSp":
            counts["connectors"] += 1
        elif tag == f"{{{NS_P_}}}grpSp":
            counts["groups"] += 1
        elif tag == f"{{{NS_P_}}}sp":
            if _has_content_text(child):
                counts["text_boxes"] += 1
            else:
                counts["shapes"] += 1
    return counts


# ------------------------------------------------------------------
# Title / divider slide detection
# ------------------------------------------------------------------

def is_title_slide(slide, slide_index):
    if slide_index == 0:
        logger.info("Slide %d — SKIP (cover/title slide)", slide_index + 1)
        return True
    try:
        texts            = []
        has_picture      = False
        has_non_ph_shape = False

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_picture = True
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip().lower())
            if not shape.is_placeholder:
                has_non_ph_shape = True

        combined   = " ".join(texts)
        word_count = len(combined.split()) if combined.strip() else 0

        for kw in ("thank you", "questions", "agenda", "table of contents"):
            if kw in combined:
                logger.info(
                    "Slide %d — SKIP (keyword match: '%s')", slide_index + 1, kw
                )
                return True

        if not has_picture and not has_non_ph_shape and word_count <= 10:
            logger.info(
                "Slide %d — SKIP (divider: %d word(s), no picture, no non-placeholder shapes)",
                slide_index + 1, word_count
            )
            return True

        return False

    except Exception as e:
        logger.error("is_title_slide error slide %d: %s", slide_index + 1, e)
        return False


# ------------------------------------------------------------------
# Build filtered + background-stripped PPTX copy
# ------------------------------------------------------------------

def create_filtered_pptx(ppt_path, skip_title_slides=True):
    """
    1. Decide which slides to keep.
    2. Copy entire PPTX to temp dir (preserves all media rels).
    3. Strip master/layout chrome + force white bg.
    4. Force white bg on kept slides (never remove slide shapes).
    5. Delete unwanted slides.
    6. Save.
    """
    source_prs = Presentation(ppt_path)
    total      = len(source_prs.slides)

    logger.info("=" * 60)
    logger.info("PPT PROCESSING STARTED")
    logger.info("  File      : %s", os.path.basename(ppt_path))
    logger.info("  Slides    : %d total", total)
    logger.info("  Skip titles: %s", skip_title_slides)
    logger.info("=" * 60)

    keep_indices = []
    skip_indices = []
    for i, slide in enumerate(source_prs.slides):
        if skip_title_slides and is_title_slide(slide, i):
            skip_indices.append(i + 1)
        else:
            keep_indices.append(i)
            summary = _slide_shape_summary(slide)
            logger.info(
                "Slide %d — KEEP | shapes=%d  pics=%d  text_boxes=%d  "
                "connectors=%d  groups=%d  other=%d",
                i + 1,
                summary["total"],
                summary["pictures"],
                summary["text_boxes"],
                summary["connectors"],
                summary["groups"],
                summary["shapes"],
            )

    logger.info("-" * 60)
    logger.info(
        "Slide filter result: %d kept, %d skipped %s",
        len(keep_indices),
        len(skip_indices),
        f"(skipped slides: {skip_indices})" if skip_indices else "",
    )
    logger.info("-" * 60)

    temp_dir     = tempfile.mkdtemp()
    working_copy = os.path.join(temp_dir, "working.pptx")
    shutil.copy2(ppt_path, working_copy)

    work_prs   = Presentation(working_copy)
    xml_slides = work_prs.slides._sldIdLst
    REL_NS     = f"{{{NS_R}}}"

    slide_rids = [
        xml_slides[i].get("r:id") or xml_slides[i].get(f"{REL_NS}id")
        for i in range(total)
    ]

    # ---- Clean ALL masters + ALL their layouts ----
    cleaned_master_ids = set()
    for slide in work_prs.slides:
        master    = slide.slide_layout.slide_master
        master_id = id(master)
        if master_id in cleaned_master_ids:
            continue
        cleaned_master_ids.add(master_id)

        logger.info("MASTER/LAYOUT CLEANUP")
        _set_white_bg(master._element, "  master")
        _strip_master_layout_chrome(master._element, "  master")

        layout_stripped_total = 0
        for li, layout in enumerate(master.slide_layouts):
            _set_white_bg(layout._element, f"  layout[{li}]")
            # Count before strip for summary
            cSld = layout._element.find(f"{{{NS_P}}}cSld")
            spTree = cSld.find(f"{{{NS_P}}}spTree") if cSld is not None else None
            before = len([c for c in (spTree or [])
                          if not _should_keep_in_master_or_layout(c)])
            _strip_master_layout_chrome(layout._element, f"  layout[{li}]")
            layout_stripped_total += before

        logger.info(
            "  Master + %d layouts cleaned | total chrome shapes removed: %d",
            len(master.slide_layouts), layout_stripped_total
        )

    # ---- Kept slides: white bg only ----
    logger.info("SLIDE BACKGROUND OVERRIDE")
    for i in keep_indices:
        try:
            _set_white_bg(work_prs.slides[i]._element, f"  slide {i + 1}")
            logger.info("  Slide %d: white background applied (all %d shapes preserved)",
                        i + 1,
                        _slide_shape_summary(work_prs.slides[i])["total"])
        except Exception as e:
            logger.warning("  Slide %d bg failed: %s", i + 1, e)

    # ---- Delete unwanted slides ----
    remove = sorted([j for j in range(total) if j not in keep_indices], reverse=True)
    if remove:
        logger.info("REMOVING %d SLIDE(S): %s", len(remove), [r + 1 for r in remove])
    for i in remove:
        rId = slide_rids[i]
        xml_slides.remove(xml_slides[i])
        try:
            work_prs.part.drop_rel(rId)
        except Exception as e:
            logger.warning("  drop_rel slide %d: %s", i + 1, e)

    out_path = os.path.join(temp_dir, "filtered.pptx")
    work_prs.save(out_path)
    logger.info("Filtered PPTX saved → %s", out_path)
    return out_path, temp_dir


# ------------------------------------------------------------------
# Main entry: PPT → filtered PPTX → PDF → PNG list
# ------------------------------------------------------------------

def render_ppt_slides_to_images(ppt_path, skip_title_slides=True):
    try:
        logger.info("render_ppt_slides_to_images called")
        logger.info("  Input file       : %s", ppt_path)
        logger.info("  Skip title slides: %s", skip_title_slides)

        filtered_path, temp_dir = create_filtered_pptx(
            ppt_path, skip_title_slides=skip_title_slides
        )

        logger.info("LIBREOFFICE: PPTX → PDF")
        result = subprocess.run(
            [LIBREOFFICE_PATH, "--headless", "--convert-to", "pdf",
             filtered_path, "--outdir", temp_dir],
            capture_output=True, text=True
        )
        if result.returncode != 0:
            logger.error("LibreOffice failed (exit %d): %s", result.returncode, result.stderr)
            raise RuntimeError(f"LibreOffice failed: {result.stderr}")

        pdf_files = [
            os.path.join(temp_dir, f)
            for f in os.listdir(temp_dir)
            if f.endswith(".pdf")
        ]
        if not pdf_files:
            raise Exception("LibreOffice produced no PDF")

        pdf_path = pdf_files[0]
        logger.info("  PDF created: %s", pdf_path)

        logger.info("PDF2IMAGE: PDF → PNG (dpi=200)")
        pages = convert_from_path(pdf_path, dpi=200, poppler_path=POPPLER_PATH)
        logger.info("  Pages extracted: %d", len(pages))

        final_images = []
        for idx, page in enumerate(pages):
            img_path = os.path.join(temp_dir, f"slide_{idx + 1}.png")
            page.save(img_path, "PNG")
            final_images.append(img_path)
            w, h = page.size
            logger.info("  PNG [%d/%d]: %s  (%dx%d px)",
                        idx + 1, len(pages), os.path.basename(img_path), w, h)

        logger.info("=" * 60)
        logger.info("PPT PROCESSING COMPLETE")
        logger.info("  Total PNG images : %d", len(final_images))
        logger.info("  Output dir       : %s", temp_dir)
        logger.info("=" * 60)
        return final_images

    except Exception as e:
        logger.error("render_ppt_slides_to_images FAILED: %s", e)
        return []
